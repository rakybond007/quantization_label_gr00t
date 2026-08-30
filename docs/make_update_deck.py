"""Build the 07/21 progress-update deck in the SAME format as
Action_quantization_0707.pdf: 16:9, white bg, bold blue left title,
black bullets (● / ○), booktabs-style tables with green-highlighted best rows,
closing TODO slide."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE = RGBColor(0x1F, 0x78, 0xC8)     # title blue (matches 0707 deck)
BLACK = RGBColor(0x11, 0x11, 0x11)
GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0xE2, 0xEF, 0xDA)    # best-row highlight (as in 0707 tables)
LINE = RGBColor(0x33, 0x33, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def title(s, text, size=27):
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.bold = True; r.font.size = Pt(size); r.font.color.rgb = BLUE
    r.font.name = "Arial"
    return tb

def bullets(s, items, top=1.25, left=0.65, width=12.0, size=15.5, height=None):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                              Inches(height if height else 5.6))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for lvl, txt, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        mark = "●  " if lvl == 0 else "○  "
        r = p.add_run(); r.text = mark + txt
        r.font.size = Pt(size if lvl == 0 else size - 1.5)
        r.font.color.rgb = BLACK
        r.font.bold = bool(bold)
        r.font.name = "Arial"
        p.space_after = Pt(7 if lvl == 0 else 4)
    return tb

def table(s, rows, col_w, top, left=0.9, hl_rows=(), font=13, header_line=True):
    n_r, n_c = len(rows), len(rows[0])
    width = Emu(int(sum(col_w) * 914400))
    tbl_shape = s.shapes.add_table(n_r, n_c, Inches(left), Inches(top),
                                   width, Inches(0.32 * n_r))
    tbl = tbl_shape.table
    tbl.first_row = False; tbl.horz_banding = False
    for j, w in enumerate(col_w):
        tbl.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.margin_left = Pt(8); c.margin_right = Pt(8)
            c.fill.solid()
            c.fill.fore_color.rgb = GREEN if i in hl_rows else RGBColor(0xFF, 0xFF, 0xFF)
            tf = c.text_frame; tf.word_wrap = False
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(font); r.font.name = "Arial"
            r.font.color.rgb = BLACK
            r.font.bold = (i == 0) or (i in hl_rows and j == 0)
            if j > 0:
                p.alignment = PP_ALIGN.CENTER
    return tbl_shape

# ---------------- Slide 1 : cover (0707 표지 포맷) ----------------
s = slide()
tb = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Action Quantization : Progress Update"
r.font.bold = True; r.font.size = Pt(34); r.font.color.rgb = BLUE; r.font.name = "Arial"
tb = s.shapes.add_textbox(Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Hojin Jeon"
r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = BLACK; r.font.name = "Arial"
tb = s.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "July 21, 2026"
r.font.size = Pt(17); r.font.color.rgb = BLACK; r.font.name = "Arial"

# ---------------- Slide 2 : Recap ----------------
s = slide(); title(s, "Recap : Where we left off (07/07)")
bullets(s, [
    (0, "Self-evolving multi-agent gate (frozen VLM judge + evolving guidance prompt) recovered the success lost by naive K=2 quantization on Robocasa & Libero.", 0),
    (1, "Gemma4-12B (bold) vs Cosmos3-nano (conservative) judge personalities confirmed on both benchmarks.", 0),
    (0, "Open TODOs from 07/07 : ① other VLM judge  ② more evolve iterations  ③ more benchmarks  ④ granularity beyond 2×  ⑤ prompt / evolver tuning", 0),
    (0, "This update : ⑤→②→③→④ addressed in order — evolver redesign, TTL call reduction, full 2×2×TTL matrix, and a causal study that closes the granularity question.", 1),
])

# ---------------- Slide 3 : Method Update (1) evolver ----------------
s = slide(); title(s, "Method Update (1) : Evolver Agent Redesign")
bullets(s, [
    (0, "Diagnosis : accept/reject rule (not the meta-prompt) caused failures — asymmetric thresholds accepted “give-up-compression” candidates and rejected genuine Pareto moves.", 0),
    (0, "Composite accept-gating (v3) : success & steps normalized to each benchmark’s measured raw↔K2 range; weighted composite score + corridor floors. Verified by replaying all recorded trajectories.", 1),
    (0, "Per-suite history : evolver now sees WHERE each edit landed (per-suite success / quant%), so it recombines the working half of a rejected candidate.", 0),
    (0, "Benchmark-context injection : generic meta-prompt + per-benchmark fact file — adding a new benchmark = one context file (enabler for TODO ③).", 0),
    (0, "New rules : bisect-after-overshoot (5–15pp quant moves near frontier), escalate-on-no-response, pair every broad YES with a HARMFUL-cluster guard.", 0),
])

# ---------------- Slide 4 : Method Update (2) TTL ----------------
s = slide(); title(s, "Method Update (2) : TTL Policy — Half the VLM Calls")
bullets(s, [
    (0, "Async 2-GPU setting hides per-call latency; the remaining cost is the NUMBER of judge calls.", 0),
    (0, "Confidence-TTL latch : reuse the previous decision while confidence is far from τ; re-judge immediately when ambiguous.", 0),
    (0, "Gripper trigger (free & predictive) : any open↔close transition inside the VLA’s own predicted chunk forces a fresh call — make-or-break moments are never latched over.", 1),
    (0, "Measured judge latency (A100, single forward P(YES) readout) :", 0),
    (1, "Gemma4-12B  175.9 / 275.4 ms (2/3 views; torch.compile → 154.9)   ·   Cosmos3-nano  79.3 / 94.8 ms", 0),
    (0, "Result : ~50% call rate on every bench×judge cell with NO quality cost (same success / steps as always-call).", 1),
])

# ---------------- Slide 5 : Results (1) Libero + TTL ----------------
s = slide(); title(s, "Experimental Results (1) - Libero w Gemma4 & Cosmos3 + TTL")
bullets(s, [
    (0, "Improved evolver + TTL, naive guidance start, 5 cycles, 40 tasks × 50 ep per cycle.", 0),
    (0, "Cosmos : evolver escapes the conservative plateau in ONE cycle (5%→26%) — new best 136 steps at 0.953, with ~60% of the calls.", 1),
    (0, "Gemma : naive guidance (47%) remains optimal; gating correctly defends it against 5 aggressive candidates.", 0),
], top=1.15, height=2.2)
table(s, [
    ["Method", "Quant.", "Calls", "Success", "Succ.-only steps"],
    ["No quantization (raw)", "–", "–", "0.932", "184"],
    ["Naive quantization (K=2)", "100%", "–", "0.800", "111"],
    ["Gemma gate + TTL (best = naive guidance)", "47%", "~50%", "0.918", "128"],
    ["Cosmos gate + TTL (best = cycle 2)", "26%", "~60%", "0.953", "136"],
], [4.6, 1.2, 1.2, 1.4, 2.0], top=3.55, hl_rows={4})

# ---------------- Slide 6 : Results (2) Robocasa + TTL ----------------
s = slide(); title(s, "Experimental Results (2) - Robocasa w Gemma4 & Cosmos3 + TTL")
bullets(s, [
    (0, "Bidirectional convergence : whichever way the judge starts wrong, the evolver walks it to the frontier.", 1),
    (1, "Gemma over-quantized at naive (69%, succ 0.615) → evolver pulled BACK to 55%, recovering 0.646.", 0),
    (1, "Cosmos under-quantized at naive (4%) → evolver pushed UP to 38% at equal success.", 0),
], top=1.15, height=2.0)
table(s, [
    ["Method", "Quant.", "Calls", "Success", "Succ.-only steps"],
    ["No quantization (raw)", "–", "–", "0.657", "330"],
    ["Naive quantization (K=2)", "100%", "–", "0.609", "221"],
    ["Gemma gate + TTL (best = cycle 2)", "55%", "~50%", "0.646", "272"],
    ["Cosmos gate + TTL (best = cycle 4)", "38%", "~60%", "0.645", "286"],
], [4.6, 1.2, 1.2, 1.4, 2.0], top=3.35, hl_rows={3, 4})

# ---------------- Slide 7 : Results (3) K3/K4 + varK ----------------
s = slide(); title(s, "Experimental Results (3) - Beyond 2× : naive K=3,4 vs bound-aware varK")
bullets(s, [
    (0, "Naive K≥3 collapses success — and even SUCCEEDING episodes gain almost no speed over K2.", 0),
    (0, "varK (magnitude-aware merge) : merge consecutive deltas only while the block sum stays inside the controller bound; never across a gripper transition.", 0),
    (0, "varK preserves success under forced high compression (Libero 0.438→0.821, above K2) — but is not faster than K2 : a success-preserving fallback, not a speed knob.", 1),
], top=1.15, height=2.15)
table(s, [
    ["", "Libero  Succ", "Steps", "Robocasa  Succ", "Steps"],
    ["Naive K=2", "0.800", "111", "0.609", "221"],
    ["Naive K=3", "0.438", "117", "0.497", "212"],
    ["Naive K=4", "0.261", "118", "0.400", "228"],
    ["varK3 (bound 0.95)", "0.821", "139", "0.626", "231"],
    ["varK4", "0.694", "145", "0.629", "238"],
], [3.2, 1.7, 1.2, 1.9, 1.2], top=3.5, hl_rows={4}, font=12.5)

# ---------------- Slide 8 : Results (4) causal ladder ----------------
s = slide(); title(s, "Experimental Results (4) - Why K2 Saturates : Causal Unlock Ladder")
bullets(s, [
    (0, "Question : with the sim’s limits removed in code, do successful episodes reach the theoretical reduction (K3=1.5×, K4=2.0× vs K2)?", 0),
    (0, "Unlocked in stairs : command clip (×K) → actuator torque caps + OSC stiffness (×K). Paired on tasks succeeding in both settings.", 0),
], top=1.15, height=1.7)
table(s, [
    ["Unlock level", "Libero K3", "Libero K4", "Robocasa K3", "Robocasa K4"],
    ["naive (as-is)", "1.04×", "1.04×", "1.15×", "1.23×"],
    ["+ command clip removed (clipK)", "0.96×", "1.00×", "1.11×", "1.21×"],
    ["+ torque & stiffness ×K (dynK)", "0.97×", "1.05×", "1.05×", "1.18×"],
    ["theoretical", "1.50×", "2.00×", "1.50×", "2.00×"],
], [3.9, 1.5, 1.5, 1.6, 1.6], top=3.0, font=12.5)
bullets(s, [
    (0, "Every unlock leaves the ratio (and success) unchanged → neither clipping nor robot dynamics binds.", 0),
    (0, "Binding constraint = the policy-controller closed loop itself (50 ms transient + replan re-anchoring). K2 IS the saturation point — so the contribution is WHERE K2 applies (our gate), not raising K.", 1),
], top=5.15, height=1.9, size=14.5)

# ---------------- Slide 9 : TODO ----------------
s = slide(); title(s, "TODO")
bullets(s, [
    (0, "Re-evolve Libero w Gemma4 with the redesigned evolver (can it now beat the naive guidance?).", 0),
    (0, "Third VLM judge : Qwen3.5 (pipeline accepts it as a model-id swap).", 0),
    (0, "Paper assembly : 10-way baseline table + 2×2×TTL matrix + K2-saturation causal study as the motivation/analysis section.", 0),
    (0, "(Optional) Robocasa w Gemma4 non-TTL rerun for full symmetry ; real-world gate demo.", 0),
])

out = "/sjw_alinlab/home/hojin2/quantization_agent_workspace/docs/Action_quantization_progress_0721.pptx"
prs.save(out)
print("saved:", out)
